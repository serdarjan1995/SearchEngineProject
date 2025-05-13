import logging

import numpy as np
import mimetypes
import requests
from bs4 import BeautifulSoup, XMLParsedAsHTMLWarning
from PIL import Image
from io import BytesIO
from pdf2image import convert_from_bytes
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import easyocr
import warnings
import urllib3

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

warnings.filterwarnings("ignore", category=XMLParsedAsHTMLWarning)


class URLScrapeProcessor:
    def __init__(self, languages=['en'], use_gpu=False):
        self.reader = easyocr.Reader(languages, gpu=use_gpu, verbose=False)
        self.logger = logging.getLogger(__name__)
        self.logger.setLevel(logging.INFO)
        ch = logging.StreamHandler()
        formatter = logging.Formatter(
            '%(asctime)s - %(name)s - %(levelname)s - %(message)s @ %(filename)s:%(funcName)s:%(lineno)d')
        ch.setFormatter(formatter)
        self.logger.addHandler(ch)

    def ocr_image(self, content):
        try:
            img = Image.open(BytesIO(content)).convert("RGB")
            result = self.reader.readtext(np.array(img), detail=0, paragraph=True)
            return "\n".join(result)
        except Exception as e:
            print(f"Image OCR failed: {e}")
            return ""

    def ocr_pdf(self, content):
        try:
            text = extract_text(BytesIO(content))
            if text.strip():
                return text
        except Exception:
            pass
        try:
            images = convert_from_bytes(content)
            result = []
            for img in images:
                img_rgb = img.convert("RGB")
                result.extend(self.reader.readtext(np.array(img_rgb), detail=0, paragraph=True))
            return "\n".join(result)
        except Exception as e:
            print(f"PDF OCR fallback failed: {e}")
            return ""

    def extract_docx_text(self, content):
        try:
            doc = Document(BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            print(f"DOCX extract failed: {e}")
            return ""

    def extract_ppt_text(self, content):
        try:
            prs = Presentation(BytesIO(content))
            return "\n".join(
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            )
        except Exception as e:
            print(f"PPTX extract failed: {e}")
            return ""

    def get_mimetype(self, url, response):
        mime = mimetypes.guess_type(url)[0]
        if not mime and 'Content-Type' in response.headers:
            mime = response.headers['Content-Type'].split(';')[0]
        return mime

    def process_url(self, url):
        try:
            headers = {
                "User-Agent": "Mozilla/5.0",
                "Accept-Language": "en-US,en;q=0.9",
            }
            response = requests.get(url, headers=headers, timeout=20, verify=False)
            if not response.ok:
                return {'url': url, 'success': False, 'info_type': 'html', 'text': ''}
            mime = self.get_mimetype(url, response)

            if mime is None:
                return {'url': url, 'success': False, 'info_type': 'unknown', 'text': ''}

            elif mime.startswith('text/html'):
                soup = BeautifulSoup(response.text, 'html5lib')
                body = soup.body

                # Remove <script>, <style>, <link> and <footer> tags from the body
                for tag in body(['script', 'style', 'link', 'nav', 'footer']):
                    tag.decompose()  # Removes the tag from the tree

                html_text = body.get_text(separator='\n', strip=True)

                return {
                    'url': url,
                    'success': bool(html_text.strip()),
                    'info_type': 'html',
                    'text': html_text
                }

            elif mime == 'application/pdf':
                text = self.ocr_pdf(response.content)
                return {'url': url, 'success': bool(text.strip()), 'info_type': 'pdf', 'text': text}

            elif mime in [
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'application/msword'
            ]:
                text = self.extract_docx_text(response.content)
                return {'url': url, 'success': bool(text.strip()), 'info_type': 'word', 'text': text}

            elif mime in [
                'application/vnd.ms-powerpoint',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            ]:
                text = self.extract_ppt_text(response.content)
                return {'url': url, 'success': bool(text.strip()), 'info_type': 'ppt', 'text': text}

            elif mime.startswith('image/'):
                text = self.ocr_image(response.content)
                return {'url': url, 'success': bool(text.strip()), 'info_type': 'image', 'text': text}

            else:
                return {'url': url, 'success': False, 'info_type': 'unsupported', 'text': ''}

        except Exception as e:
            print(f"Error processing {url}: {e}")
            return {'url': url, 'success': False, 'info_type': 'error', 'text': ''}

    def process_urls(self, url_list):
        return [self.process_url(url) for url in url_list]


url_scrape_processor = URLScrapeProcessor()
